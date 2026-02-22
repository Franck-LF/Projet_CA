"""
    Class Document
    Class Extraction
    Class ParamExtraction

    ...
"""




import os
import logging
from typing import List, Dict
import hashlib
import pandas as pd
import pdfplumber
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
import xlrd # To read xls files




class MyDocument():
    ''' Objet document '''

    def __init__(self, 
                 title:str = "",
                 path:str = "",
                 text:str = "",
                 author:str = "",
                 date:str = "",
                 extension:str = "",
                 unit:str = "",
                 hash:str = "",
                 has_changed:bool = False,
                 meta:Dict = {}):
        ''' Document initialization
            title: document title,
            path: document path,
            text: full document content,
            author: document author,
            date: document last modification date,
            extension: document extension,
            unit: ABP unit (CAR, ACAP, UGA, USP ...)
            hash: string with the file hash (hashkey: sha256)
                  used to check if a file has been modified,
            has_changed: boolean to indicate if document has changed (according to previous document process)
                         if False we do not need to extract the document text.
            meta: UNUSED.
        '''
        self.title:str = title
        self.path:str = path
        self.text:str = text
        self.author:str = author
        self.date:str = date
        self.extension:str = extension
        self.unit:str = unit
        self.hash:str = hash
        self.has_changed:bool = has_changed
        self.meta:Dict = meta

    def __str__(self) -> str:
        ''' Renvoie une chaîne de caractères avec les informations de l'objet 'MyDocument' '''
        st = f"-------------------\nDocument : {self.title}.{self.extension}"
        assert self.path
        st += f"\nPath: {self.path}"
        if self.author: st += f"\nAuthor: {self.author}"
        if self.date: st += f"\nDate: {self.date}"
        assert self.unit
        st += f"\nUnit : {self.unit}"
        st += f"\nHas changed : {self.has_changed}"
        return st

    def print_document(self, verbose = 1):
        print('-' * 10, self, '-' * 10)
        if verbose == 1:
            print(self.text)


class ParamExtraction():
    ''' Paramètres de l'extraction '''
    path:str = ""
    real_path:str = ""
    path_temp_pdf:str = ""
    max_files:int = 1000
    recurrent:bool = True
    extract_text:bool = False
    convert_to_pdf:bool = False
    verbose:int = 0
    lst_filetypes:List[str] = []
    logger:logging.Logger = None


class Extraction():
    ''' Objet extraction '''

    def __init__(self, param:ParamExtraction, df_old:pd.DataFrame = None):
        ''' Initialise l'objet Extraction '''
        self._list_documents:List[MyDocument] = []
        self.set_param(param)
        self._list_not_converted:List = []
        self._list_converted:List = []
        self._list_xlsm_problems:List = []
        self._df_old:pd.DataFrame = df_old

    def __str__(self) -> str:
        ''' Renvoie une chaîne de caractères avec les informations de l'objet 'Extraction' '''
        st = f"Extraction - nombre de documents  {len(self._list_documents)}\n"
        st += '\n'.join([doc.__str__() for doc in self._list_documents])
        return st

    def __len__(self) -> int:
        ''' Renvoie le nombre de documents '''
        return len(self._list_documents)

    def print_list_files(self):
        if self._list_documents:
            for document in self._list_documents:
                print(f'Document : "{document.path}{document.title}.{document.extension}"')
        else:
            print(f'Pas de fichier analysés dans le répertoire "{self.path}"')

    def print_list_not_converted(self):
        if self._list_not_converted:
            for document in self._list_not_converted:
                print(f'Document : "{document.path}{document.title}.{document.extension}"')

    def print_list_converted(self):
        if self._list_converted:
            for document in self._list_converted:
                print(f'Document : "{document.path}{document.title}.{document.extension}"')

    def print_files_stats(self) -> None:
        print(f"Nombre de pdf  : {len([True for doc in self._list_documents if doc.extension == 'pdf'])}")
        print(f"Nombre de pptx : {len([True for doc in self._list_documents if doc.extension == 'pptx'])}")
        print(f"Nombre de docx : {len([True for doc in self._list_documents if doc.extension == 'docx'])}")
        print(f"Nombre de doc : {len([True for doc in self._list_documents if doc.extension == 'doc'])}")
        print(f"Nombre de xlsx : {len([True for doc in self._list_documents if doc.extension == 'xlsx'])}")
        print(f"Nombre de xlsm : {len([True for doc in self._list_documents if doc.extension == 'xlsm'])}")
        print(f"Nombre de xls : {len([True for doc in self._list_documents if doc.extension == 'xls'])}")
        print(f"Total : {len(self)}")
        self._param.logger.info('-' * 40)
        self._param.logger.info(" Statistiques d'extraction :")
        self._param.logger.info(f" Nombre de pdf  : {len([True for doc in self._list_documents if doc.extension == 'pdf'])}")
        self._param.logger.info(f" Nombre de pptx : {len([True for doc in self._list_documents if doc.extension == 'pptx'])}")
        self._param.logger.info(f" Nombre de docx : {len([True for doc in self._list_documents if doc.extension == 'docx'])}")
        self._param.logger.info(f" Nombre de doc : {len([True for doc in self._list_documents if doc.extension == 'doc'])}")
        self._param.logger.info(f" Nombre de xlsx : {len([True for doc in self._list_documents if doc.extension == 'xlsx'])}")
        self._param.logger.info(f" Nombre de xlsm : {len([True for doc in self._list_documents if doc.extension == 'xlsm'])}")
        self._param.logger.info(f" Nombre de xls : {len([True for doc in self._list_documents if doc.extension == 'xls'])}")
        self._param.logger.info(f" Total : {len(self)}")
        self._param.logger.info('-' * 40)
        # print(f"Nb de documents convertis en PDF : {len(self._list_converted)}")
        # print(f"Nb de documents non convertis en PDF (car PB) : {len(self._list_not_converted)}")

    def set_param(self, param:ParamExtraction):
        param.lst_filetypes = [item.lower() for item in param.lst_filetypes]
        self._param = param

    def reset_param(self):
        self._param = None

    def to_dataframe(self):
        return pd.DataFrame.from_dict({
            'Extensions': [doc.extension for doc in self._list_documents],
            'Titles': [doc.title for doc in self._list_documents],
            'Texts': [doc.text for doc in self._list_documents],
            'Author': [doc.author for doc in self._list_documents],
            'Date': [doc.date for doc in self._list_documents],
            'Path': [doc.path for doc in self._list_documents],
            'Unit': [doc.unit for doc in self._list_documents],
            'Hash': [doc.hash for doc in self._list_documents],
            'Has_changed': [doc.has_changed for doc in self._list_documents],
            })


    def create_documents_from_directory(self) -> bool:
        ''' Scan directory and sub directories
            Create document for each file
        '''
        # Create temporary file to put pdf (Needs to be scanned before the process)
        # if not self.create_temp_directory():
        #     return False

        if not all([item in os.listdir(self._param.path) for item in ['ACAP', 'CAR', 'UGA', 'USP']]):
            print("Erreur : L'architecture des sous-dossiers d'unités est incorrecte")
            self._param.logger.error("L'architecture des sous-dossiers d'unités est incorrecte")
            return False

        # Scan root folder
        return self.create_documents_from_directory_REC(self._param.path, 0)


    def create_documents_from_directory_REC(self, path:str, depth:int, unit:str = '') -> bool:
        ''' Recursive function to scan directory and sub directories and create document for each file '''

        assert path[-1] == "/" or path[-1] == "\\"
        for file in os.listdir(path):
            # Scan folder
            if os.path.isdir(path + file):
                if len(os.listdir(path + file)) == 0: continue
                if depth == 0:
                    # First Subfolder (depth 0), we add it as "UNIT" to the document
                    unit = file
                    if self._param.verbose == 1:
                        print("UNIT: ", file)

                if self._param.recurrent:
                    # Scan subfolder
                    if not self.create_documents_from_directory_REC(path + file + '\\', depth + 1, unit):
                        return False

            elif depth > 0:
                if file[0] == '~': continue
                if self._param.verbose == 1:
                    print(file)

                # add document to the list
                document = self.create_document(path, file, unit)

                if document != None:
                    self._list_documents.append(document)

                    if self._param.verbose == 1:
                        print(f"Title: {document.title}")
                        if document.author != '':
                            print(f"Author: {document.author}")
                        if document.date != '':
                            print(f"Date: {document.date}")
                        print(f"Path: {path}")
                        # extension, title, text, author, date, meta = "", "", "", "", "", {}

                    if len(self._list_documents) == self._param.max_files:
                        return False
                else:
                    print(f"Impossible de créer le document : {path + file}")

        return True

    def create_document(self, path:str, filename:str, unit:str) -> MyDocument:
        ''' Create document and add it to the list.
            (extract informations: text, author, date, meta ...)

            Arg:
            - path: full path of the file (C:/Documents/../).
            - filename: filename with extension (filename.pdf).
        '''

        @staticmethod
        def format_author(author:str) -> str:
            if '.' in author:
                author = author.replace('.', ' ')
            items = author.lower().split(' ')
            return ' '.join(item.strip()[0].upper() + item.strip()[1:] for item in items if item.strip() != '')

        @staticmethod
        def extract_text_and_metadata_from_pdf(path_and_filename:str) -> str:
            ''' Extract text and metadata from PDF file.

            Return: 
            - text:  full text extracted from the file,
            - author: author of the document (if exists),
            - modDate: last modified date.

            Arg:
            - path_and_filename : full path + name of file + extension (C:/Documents/..../filename.pdf).
            '''
            text, author, modDate = "", "", ""
            with pdfplumber.open(path_and_filename) as pdf:
                for i, page in enumerate(pdf.pages):
                    current_text = page.extract_text()
                    text += ' ' + current_text.strip()

                if "Author" in pdf.metadata.keys():
                    author = format_author(pdf.metadata["Author"])
                if "ModDate" in pdf.metadata.keys():
                    modDate = pdf.metadata["ModDate"]
                    modDate = f"{modDate[8:10]}/{modDate[6:8]}/{modDate[2:6]}"
                elif 'CreationDate' in pdf.metadata.keys():
                    modDate = pdf.metadata["CreationDate"]
                    modDate = f"{modDate[2:6]}-{modDate[6:8]}-{modDate[8:10]}"

            return text, author, modDate
            
        @staticmethod
        def extract_text_and_metadata_from_pptx(path_and_filename) -> str:
            ''' Extract text and metadata from Powerpoint file.
            
            Return: 
            - text:  full text extracted from the file,
            - author: author of the document (if exists),
            - modDate: last modified date.

            Arg:
            - path_and_filename : full path + name of file + extension ('C:/Documents/..../filename.pptx).
            '''
            text, author, modDate = "", "", ""
            prs = Presentation(path_and_filename)
            if prs:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += ' ' + run.text.strip()
                prop = prs.core_properties
                author = format_author(prop.author)
                modDate = prop.modified.strftime("%d/%m/%Y")

            return text, author, modDate
        
        @staticmethod
        def extract_text_and_metadata_from_docx(path_and_filename) -> str:
            ''' Extract text from Word file.

            Return:
            - text:  full text extracted from the file,
            - author: author of the document (if exists),
            - modDate: last modified date.

            Arg:
            - path_and_filename : full path + name of file + extension ('C:/Documents/..../filename.docx).
            '''
            text, author, modDate = "", "", ""
            doc = Document(path_and_filename)
            if doc:
                for paragraph in doc.paragraphs:
                    text += '\n' + paragraph.text.strip()
                prop = doc.core_properties
                author = format_author(prop.author)
                modDate = prop.modified.strftime("%d/%m/%Y")
            return text, author, modDate
        
        @staticmethod
        def extract_text_and_metadata_from_doc(path_and_filename) -> str:
            ''' Extract text from doc file (old Word files).

            Return:
            - title: title of the file,
            - text:  full text extracted from the file.

            Arg:
            - path_and_filename : full path + name of file + extension ('C:/Documents/..../filename.doc).
            '''
            text, author, modDate = "", "", ""
            # TODO
            assert False
            return text, author, modDate

        @staticmethod
        def extract_text_and_metadata_from_xlsx_or_xlsm(path_and_filename) -> str:
            ''' Extract text from xlsx file.

            Return:
            - title: title of the file,
            - text:  full text extracted from the file.

            Arg:
            - path_and_filename : full path + name of file + extension ('C:/Documents/..../filename.xlsx).
            '''
            text, author, modDate = "", "", ""
            try:
                workbook = load_workbook(path_and_filename, read_only = True)
                if workbook:
                    for worksheet in workbook:
                        for row in worksheet.iter_rows():
                            for cell in row:
                                if(cell.value != None):
                                    text += ' ' + str(cell.value).strip()
                    # TODO
                    # author = ...
                    # modDate = ...
            except:
                self._list_xlsm_problems.append(path + filename)
                if self._param.verbose == 1:
                    print(f"Cannot extract text from file: {path + filename}")

            return text, author, modDate

        @staticmethod
        def extract_text_and_metadata_from_xls(path_and_filename) -> str:
            ''' Extract text from xsl file (old Excel files).

            Return:
            - title: title of the file,
            - text:  full text extracted from the file.

            Arg:
            - path_and_filename : full path + name of file + extension ('C:/Documents/..../filename.xls).
            '''
            text, author, modDate = "", "", ""
            book = xlrd.open_workbook(path_and_filename)
            if book:
                sh = book.sheet_by_index(0)
                for rx in range(sh.nrows):
                    for cx in range(sh.ncols):
                        text += ' ' + str(sh.cell_value(rowx=rx, colx=cx)).strip()
                # TODO
                # author = ...
                # modDate = ...
            return text, author, modDate


        assert filename[0] != '~' and unit != ''
        extension, title, text, author = "", "", "", ""
        date, hash, has_changed, meta = "", "", True, {}

        try:
            with open(path + filename, "rb") as f:
                hash = hashlib.sha256(f.read()).hexdigest()
        except:
            print(F"Hash Issue: {filename}")

        bToExtract = True

        if isinstance(self._df_old, pd.DataFrame):
            # Search the hash in the old pd.DataframeCheck
            assert 'Hash' in self._df_old.columns
            df_res = self._df_old[self._df_old['Hash'] == hash]

            if not df_res.empty:
                # Hash was found in the old DataFrame
                # We retrieve the informations from it.
                if self._param.verbose == 1:
                    print('Hash found')
                # assert df_res.shape[0] == 1
                has_changed = False
                text = df_res.loc[df_res.index[0]]['Texts']
                author = df_res.loc[df_res.index[0]]['Author']
                date = df_res.loc[df_res.index[0]]['Date']
                bToExtract = False
            elif self._param.verbose == 1:
                print('Hash not found')

        if filename[-3:].lower() == 'pdf' and 'pdf' in self._param.lst_filetypes:
            extension = 'pdf'
            title = filename[:-4]
            if bToExtract and self._param.extract_text:
                text, author, date = extract_text_and_metadata_from_pdf(path + filename)

        elif filename[-4:].lower() == 'pptx' and 'pptx' in self._param.lst_filetypes:
            extension = 'pptx'
            title = filename[:-5]
            if bToExtract and self._param.extract_text:
                text, author, date = extract_text_and_metadata_from_pptx(path + filename)

        elif filename[-4:].lower() == 'docx' and 'docx' in self._param.lst_filetypes:
            extension = 'docx'
            title = filename[:-5]
            if bToExtract and self._param.extract_text:
                text, author, date = extract_text_and_metadata_from_docx(path + filename)

        elif filename[-3:].lower() == 'doc' and 'doc' in self._param.lst_filetypes:
            extension = 'doc'
            title = filename[:-4]
            if bToExtract and self._param.extract_text:
                text, author, date = extract_text_and_metadata_from_doc(path + filename)

        elif filename[-4:].lower() in ['xlsx', 'xlsm'] and filename[-4:].lower() in self._param.lst_filetypes:
            extension = filename[-4:]
            title = filename[:-5]
            if bToExtract and self._param.extract_text:
                text, author, date = extract_text_and_metadata_from_xlsx_or_xlsm(path + filename)

        elif filename[-3:].lower() == 'xls' and 'xls' in self._param.lst_filetypes:
            extension = 'xls'
            title = filename[:-4]
            if bToExtract and self._param.extract_text:
                text, author, date = extract_text_and_metadata_from_xls(path + filename)

        else:
            if self._param.verbose == 2:
                print(f"Autre extension : {filename}")
            return None

        if text.strip() == '':
            text = title

        return MyDocument(title, path, text, author, date, extension, unit, hash, has_changed, meta)


